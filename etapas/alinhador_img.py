import streamlit as st
from pptx import Presentation
import io
import pdfplumber
import google.generativeai as genai
import os
from PIL import Image

# Configuração do Gemini API
gemini_api_key = os.getenv("GEMINI_API_KEY")
genai.configure(api_key=gemini_api_key)

# Inicializa os modelos do Gemini
modelo_vision = genai.GenerativeModel("gemini-2.0-flash", generation_config={"temperature": 0.1})
modelo_texto = genai.GenerativeModel("gemini-1.5-flash")

def extract_text_from_pptx(file):
    prs = Presentation(file)
    slides_text = []
    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):  # Verifica se o shape tem texto
                slide_text += shape.text + "\n"
        slides_text.append(slide_text.strip())
    return slides_text

def extract_text_from_pdf(file):
    pages_text = []
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or "(Página sem texto extraível)"
            pages_text.append(text)
    return pages_text

st.set_page_config(layout="wide")
st.title("Upload e Extração de Texto de PPTX e PDF")

uploaded_file = st.file_uploader("Envie um arquivo .pptx, .pdf ou imagem", type=["pptx", "pdf", "jpg", "jpeg", "png"])

if uploaded_file is not None:
    file_type = uploaded_file.name.split(".")[-1]
    branding_material = ""
    
    if file_type == "pptx":
        texts = extract_text_from_pptx(io.BytesIO(uploaded_file.read()))
        branding_material = "\n\n".join(texts)
    elif file_type == "pdf":
        texts = extract_text_from_pdf(io.BytesIO(uploaded_file.read()))
        branding_material = "\n\n".join(texts)
    
    if file_type in ["pptx", "pdf"]:
        for i, text in enumerate(texts):
            st.subheader(f"Página {i+1}")
            st.text_area(f"Texto da Página {i+1}", text, height=200)
        
        # Salvar em variável global
        st.session_state["extracted_texts"] = texts
        st.success("Texto extraído com sucesso!")
    
    if file_type in ["jpg", "jpeg", "png"]:
        st.image(uploaded_file, caption='Imagem Carregada', use_column_width=True)
        image = Image.open(uploaded_file)
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format=image.format)
        img_bytes = img_byte_arr.getvalue()
        mime_type = "image/png" if image.format == "PNG" else "image/jpeg"
        
        prompt = f"""
        Você está aqui para aprovar imagens de criativos para campanhas de marketing digital para a cooperativa Holambra.
        Se atente ao mínimo e extremo detalhe de tudo que está na imagem, pois você é extremamente detalhista.
        
        O cliente Holambra já deu alguns feedbacks sobre criativos no passado.
        
        Considerando os materiais de branding do cliente e as diretrizes já existentes ({branding_material}),
        diga se a imagem seria aprovada ou não e o que precisa melhorar para ser aprovada.
        """
        
        try:
            with st.spinner('Analisando a imagem...'):
                resposta = modelo_vision.generate_content(
                    contents=[prompt, {"mime_type": mime_type, "data": img_bytes}]
                )
                descricao = resposta.text
                st.subheader('Aprovação da Imagem')
                st.write(descricao)
        except Exception as e:
            st.error(f"Ocorreu um erro ao processar a imagem: {e}")
        
    if st.button("Remover Arquivo"):
        st.session_state.clear()
        st.experimental_rerun()
